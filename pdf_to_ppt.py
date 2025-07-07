import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
import os
import tkinter as tk
from tkinter import filedialog
from tkinter import ttk
from tkinter import font
import threading
import time

def pdf_to_ppt_conversion(pdf_path, ppt_path, progress_callback, status_callback):
    """
    Convierte un PDF a un PPT, con retroalimentación de progreso y estado.

    Args:
        pdf_path (str): Ruta al archivo PDF de entrada.
        ppt_path (str): Ruta al archivo PPT de salida.
        progress_callback (callable): Función para informar el progreso (0-100).
        status_callback (callable): Función para informar el estado del proceso (texto).
    """
    try:
        status_callback("Abriendo archivo PDF...")
        pdf_document = fitz.open(pdf_path)
        total_pages = len(pdf_document)

        # Crear una nueva presentación
        status_callback("Creando presentación de PowerPoint...")
        prs = Presentation()
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)

        for page_number in range(total_pages):
            status_callback(f"Procesando página {page_number + 1} de {total_pages}...")
            page = pdf_document.load_page(page_number)
            pix = page.get_pixmap()
            image_filename = f"page_{page_number}.png"
            pix.save(image_filename)

            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            left = top = Inches(0)
            pic = slide.shapes.add_picture(image_filename, left, top, width=prs.slide_width, height=prs.slide_height)
            os.remove(image_filename)  # Limpiar archivo temporal

            # Calcular y reportar el progreso (0-100)
            progress = int((page_number + 1) / total_pages * 100)
            progress_callback(progress)

        status_callback("Guardando presentación...")
        prs.save(ppt_path)
        pdf_document.close()
        status_callback(f"Conversión exitosa. PPT guardado en: {ppt_path}")
    except Exception as e:
        status_callback(f"Ocurrió un error: {e}")
        progress_callback(0)  # Reset progress on error
    finally:
        # Asegurar que el estado final se establezca
        status_callback("Proceso completado.") # Esto asegura que el estado final se muestre
        progress_callback(100)

def start_conversion_thread(pdf_path, ppt_path, progress_var, status_var, root):
    """
    Inicia la conversión en un hilo separado para no bloquear la interfaz de usuario.

    Args:
        pdf_path (str): Ruta del PDF.
        ppt_path (str): Ruta del PPTX.
        progress_var (tk.IntVar): Variable para actualizar la barra de progreso.
        status_var (tk.StringVar): Variable para actualizar el texto de estado.
        root (tk.Tk): La ventana principal de la aplicación.
    """
    # Deshabilitar el botón de carga durante la conversión
    root.select_file_button.config(state=tk.DISABLED)

    def progress_callback(progress):
        """Actualiza la barra de progreso."""
        progress_var.set(progress)
        root.update_idletasks()  # Forzar actualización de la GUI

    def status_callback(status):
        """Actualiza el texto de estado."""
        status_var.set(status)
        root.update_idletasks()  # Forzar actualización de la GUI

    # Crear un nuevo hilo para la conversión
    conversion_thread = threading.Thread(target=pdf_to_ppt_conversion, args=(pdf_path, ppt_path, progress_callback, status_callback))
    conversion_thread.start()

    # Habilitar el botón de carga cuando el hilo termina
    def enable_button_after_completion():
        conversion_thread.join()  # Esperar a que el hilo termine
        root.select_file_button.config(state=tk.NORMAL)
    # Crear un hilo para habilitar el botón después de que termine la conversion
    enable_button_thread = threading.Thread(target=enable_button_after_completion)
    enable_button_thread.start()

def main():
    """
    Función principal que crea la interfaz gráfica y maneja la lógica de la aplicación.
    """
    root = tk.Tk()
    root.title("PDF a PPTX Converter")
    root.geometry("600x200")  # Tamaño inicial de la ventana
    root.resizable(False, False)  # Evitar que se redimensione

    # --- Estilo ---
    # Fuente personalizada
    font_style = font.Font(family="Segoe UI", size=10)
    # Colores
    bg_color = "#f0f0f0"  # Gris claro
    button_color = "#4CAF50"  # Verde
    button_hover_color = "#45a049"  # Verde más oscuro
    text_color = "#333333"  # Gris oscuro

    root.configure(bg=bg_color)

    # --- Widgets ---
    # Frame principal
    main_frame = ttk.Frame(root, padding="20", style="MainFrame.TFrame")
    main_frame.pack(fill=tk.BOTH, expand=True)
    main_frame.configure(style="MainFrame.TFrame")

     # Estilo para el Frame
    root.style = ttk.Style()
    root.style.configure("MainFrame.TFrame", background=bg_color)

    # Botón "Seleccionar archivo"
    select_file_button = tk.Button(
        main_frame,
        text="Seleccionar archivo PDF",
        command=lambda: select_file(root, progress_var, status_var),
        font=font_style,
        bg=button_color,
        fg="white",
        relief=tk.RAISED,
        borderwidth=2,
        cursor="hand2",
        activebackground=button_hover_color,
        activeforeground="white",
    )
    select_file_button.pack(pady=10, fill=tk.X)
    root.select_file_button = select_file_button  # Guardar referencia para habilitar/deshabilitar

    # Barra de progreso
    progress_var = tk.IntVar()
    progress_bar = ttk.Progressbar(
        main_frame,
        variable=progress_var,
        maximum=100,
        mode="determinate",
    )
    progress_bar.pack(pady=10, fill=tk.X)

    # Etiqueta de estado
    status_var = tk.StringVar()
    status_label = tk.Label(
        main_frame,
        textvariable=status_var,
        font=font_style,
        fg=text_color,
        bg=bg_color,
        anchor=tk.CENTER,
    )
    status_label.pack(pady=10, fill=tk.X)

    # Botón de cerrar
    close_button = tk.Button(
        root,
        text="Cerrar",
        command=root.destroy,
        font=font_style,
        bg=button_color,
        fg="white",
        relief=tk.RAISED,
        borderwidth=2,
        cursor="hand2",
        activebackground=button_hover_color,
        activeforeground="white",
    )
    close_button.place(x=520, y=160)  # Posicionamiento absoluto

    # --- Funciones ---

    def select_file(root, progress_var, status_var):
        """
        Abre el diálogo de selección de archivo y comienza la conversión.
        """
        pdf_path = filedialog.askopenfilename(
            title="Selecciona el archivo PDF",
            filetypes=[("Archivos PDF", "*.pdf")]
        )
        if pdf_path:
            # Determinar la ruta de salida del PPTX
            base_name = os.path.splitext(os.path.basename(pdf_path))[0]
            ppt_path = os.path.join(os.path.dirname(pdf_path), f"{base_name}.pptx")
            # Iniciar la conversión en un hilo separado
            start_conversion_thread(pdf_path, ppt_path, progress_var, status_var, root)

    # --- Iniciar la aplicación ---
    root.mainloop()

if __name__ == "__main__":
    main()

